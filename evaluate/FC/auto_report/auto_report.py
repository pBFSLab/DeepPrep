import os
from pathlib import Path
import pptx
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor


def add_picture(slide, pic_file, left, top, width, height):
    slide.shapes.add_picture(pic_file, Mm(left), Mm(top), Mm(width), Mm(height))
    return slide


def add_text(slide, text, left, top, width, height, font_name='Arial', font_size=18, bold=False, font_color=(0, 0, 0)):
    # 在指定位置添加文本框
    textbox = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(height))
    tf = textbox.text_frame

    para = tf.paragraphs[0]
    para.text = text

    font = para.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = RGBColor(*font_color)
    return slide


def add_vol_fc_report(prs, data_path, pipeline1, pipeline2, subj, fc_file_name):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide = add_text(slide, subj, 0, 0, 150, 10.23, bold=True)
    slide = add_text(slide, fc_file_name[3:-4], 0, 10.23, 150, 10.23, bold=True)
    slide = add_text(slide, pipeline1, 30, 43, 50, 10.23, bold=True)
    slide = add_text(slide, pipeline2, 30, 138, 50, 10.23, bold=True)
    pipeline1_fc_file = data_path / pipeline1 / subj / fc_file_name
    slide = add_picture(slide, str(pipeline1_fc_file), 120, 0, 190, 95)
    pipeline2_fc_file = data_path / pipeline2 / subj / fc_file_name
    slide = add_picture(slide, str(pipeline2_fc_file), 120, 95.5, 190, 95)


def add_surf_fc_report(prs, data_path, pipeline1, pipeline2, subj, fc_file_name):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide = add_text(slide, subj, 0, 0, 150, 10.23, bold=True)
    slide = add_text(slide, fc_file_name[3:-5], 0, 10.23, 50, 10.23, bold=True)
    slide = add_text(slide, pipeline1, 30, 43, 50, 10.23, bold=True)
    slide = add_text(slide, pipeline2, 30, 138, 50, 10.23, bold=True)
    pipeline1_fc_file = data_path / pipeline1 / subj / fc_file_name
    slide = add_picture(slide, str(pipeline1_fc_file), 101.26, 0, 136.15, 95)
    pipeline2_fc_file = data_path / pipeline2 / subj / fc_file_name
    slide = add_picture(slide, str(pipeline2_fc_file), 101.26, 95.5, 136.15, 95)


def main(pipeline1, pipeline2, dist_dir, dataset, ppt_outpath):
    data_path = Path(f'{dist_dir}/{dataset}/derivatives/analysis')
    subjs = os.listdir(data_path / pipeline1)
    subjs = sorted(subjs)
    prs = pptx.Presentation('report_template.pptx')
    for idx, subj in enumerate(subjs):
        print(f'{idx + 1}/{len(subjs)}: {subj}')
        vol_fc_files = sorted((data_path / pipeline2 / subj).glob('*.png'))
        for vol_fc_file in vol_fc_files:
            add_vol_fc_report(prs, data_path, pipeline1, pipeline2, subj, vol_fc_file.name)

        surf_fc_files = sorted((data_path / pipeline2 / subj).glob('*.tiff'))
        for surf_fc_file in surf_fc_files:
            add_surf_fc_report(prs, data_path, pipeline1, pipeline2, subj, surf_fc_file.name)
    prs.save(ppt_outpath)


if __name__ == '__main__':
    pipeline1 = 'DeepPrep'
    pipeline2 = 'DeepPrep-SDC'
    bids_dir = '/mnt/ngshare/DeepPrep'
    dataset = 'MSC'
    save_path = Path(bids_dir) / dataset / 'derivatives' / 'analysis'
    ppt_outpath = save_path / f'{dataset}_{pipeline1}_{pipeline2}_new.pptx'
    main(pipeline1, pipeline2, bids_dir, dataset, ppt_outpath)
